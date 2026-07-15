"""traffic_law_z.ipynb 소개용 발표자료(PPT)와 다이어그램 생성 스크립트.

생성물:
- traffic_law_z_workflow.png : 전체 상담 워크플로 흐름도
- traffic_law_z_langgraph.png : LangGraph node/edge/state 다이어그램
- traffic_law_z_overview.pptx : 발표자료 초안

실행:
    python docs/build_traffic_law_z_ppt.py
"""
from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).parent
TEST_DATASET_DIR = DOCS.parent / "test_dataset"
EVAL_CSV = TEST_DATASET_DIR / "evaluation_results_9b5c2dbe.csv"
EVAL_WO_BM25_CSV = TEST_DATASET_DIR / "evaluation_results_a5cb495c.csv"
WORKFLOW_IMG = DOCS / "traffic_law_z_workflow.png"
LANGGRAPH_IMG = DOCS / "traffic_law_z_langgraph.png"
EVAL_SUMMARY_IMG = DOCS / "traffic_law_z_eval_summary.png"
EVAL_WO_BM25_SUMMARY_IMG = DOCS / "traffic_law_z_eval_summary_wo_bm25.png"
OUT_PPTX = DOCS / "traffic_law_z_overview.pptx"

# ── 테마 색상 ────────────────────────────────────────────────
DARK = RGBColor(0x14, 0x2A, 0x4F)      # 딥 네이비
ACCENT = RGBColor(0x1F, 0x77, 0xB4)    # 블루
ACCENT2 = RGBColor(0x2C, 0xA0, 0x2C)   # 그린
WARN = RGBColor(0xC0, 0x39, 0x2B)      # 레드
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# matplotlib 한글 폰트
for _name in ("Malgun Gothic", "AppleGothic", "NanumGothic"):
    if any(f.name == _name for f in fm.fontManager.ttflist):
        plt.rcParams["font.family"] = _name
        break
plt.rcParams["axes.unicode_minus"] = False

MPL_DARK = "#142A4F"
MPL_ACCENT = "#1F77B4"
MPL_GREEN = "#2CA02C"
MPL_ORANGE = "#E08A1E"
MPL_GRAY = "#5B6B7B"


# ── 다이어그램 공통 유틸 ─────────────────────────────────────
def _box(ax, xy, w, h, text, face, edge, fontsize=11, fontcolor="white", radius=0.03):
    x, y = xy
    box = FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0.01,rounding_size={radius}",
        linewidth=1.6, edgecolor=edge, facecolor=face, zorder=2,
    )
    ax.add_patch(box)
    ax.text(
        x + w / 2, y + h / 2, text,
        ha="center", va="center", fontsize=fontsize, color=fontcolor,
        weight="bold", zorder=3, wrap=True,
    )
    return (x + w / 2, y + h / 2)


def _arrow(ax, start, end, color=MPL_GRAY, style="-|>", lw=1.8, ls="-"):
    ax.add_patch(FancyArrowPatch(
        start, end, arrowstyle=style, mutation_scale=16,
        color=color, lw=lw, linestyle=ls, zorder=1,
        shrinkA=6, shrinkB=6,
    ))


def build_workflow_diagram():
    """전체 상담 워크플로 흐름도."""
    fig, ax = plt.subplots(figsize=(12, 6.2))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6.2)
    ax.axis("off")

    # 상단 파이프라인
    _box(ax, (0.3, 4.7), 2.0, 1.0, "사용자 질문\n(교통 상담)", "#E8EEF6", MPL_ACCENT, 11, MPL_DARK)
    _box(ax, (2.9, 4.7), 2.3, 1.0, "① 쟁점 분석\n(LLM 구조화 추출)", MPL_ACCENT, MPL_ACCENT, 11)
    _box(ax, (5.8, 4.7), 2.4, 1.0, "② 이중 하이브리드\n검색 (법령·판례)", MPL_GREEN, MPL_GREEN, 11)
    _box(ax, (8.8, 4.7), 2.9, 1.0, "③ 근거 제한 답변\n(출처 인용 + 안전장치)", MPL_DARK, MPL_DARK, 11)

    _arrow(ax, (2.3, 5.2), (2.9, 5.2))
    _arrow(ax, (5.2, 5.2), (5.8, 5.2))
    _arrow(ax, (8.2, 5.2), (8.8, 5.2))

    # 쟁점 분석 산출물
    _box(ax, (2.9, 3.15), 2.3, 1.1,
         "category / issues\nsearch_queries\nlikely_articles\nmissing_facts",
         "#F2F5F9", MPL_ACCENT, 9, MPL_DARK)
    _arrow(ax, (4.05, 4.7), (4.05, 4.25), MPL_ACCENT)

    # 검색 세부 (벡터 + BM25 → RRF)
    _box(ax, (5.55, 2.75), 1.35, 0.8, "벡터 검색\n(Chroma)", "#DDEBD8", MPL_GREEN, 9, MPL_DARK)
    _box(ax, (7.05, 2.75), 1.35, 0.8, "BM25\n키워드 검색", "#DDEBD8", MPL_GREEN, 9, MPL_DARK)
    _box(ax, (6.15, 1.35), 1.9, 0.8, "RRF 순위 융합\n법령7 · 판례5", MPL_GREEN, MPL_GREEN, 9)
    _arrow(ax, (7.0, 4.7), (6.6, 3.55), MPL_GREEN)
    _arrow(ax, (7.0, 4.7), (7.6, 3.55), MPL_GREEN)
    _arrow(ax, (6.2, 2.75), (6.9, 2.15), MPL_GREEN)
    _arrow(ax, (7.7, 2.75), (7.1, 2.15), MPL_GREEN)
    _arrow(ax, (7.1, 1.75), (9.5, 4.7), MPL_GREEN, ls="--")

    # 답변 구성
    _box(ax, (8.8, 2.6), 2.9, 1.7,
         "잠정 판단\n→ 적용 근거[L·C]\n→ 추가 필요 사실\n→ 권장 조치\n+ 법률자문 아님 고지",
         "#E8EEF6", MPL_DARK, 9, MPL_DARK)
    _arrow(ax, (10.25, 4.7), (10.25, 4.3), MPL_DARK)

    # 메모리 루프
    _box(ax, (0.3, 2.6), 2.0, 1.0, "대화 메모리\n(thread_id)", "#FBEEDD", MPL_ORANGE, 10, "#7A4B12")
    _arrow(ax, (10.25, 2.6), (1.3, 2.6), MPL_ORANGE, ls="--")
    _arrow(ax, (1.3, 3.6), (3.5, 4.7), MPL_ORANGE, ls="--")
    ax.text(5.7, 2.42, "후속 질문 문맥 유지", fontsize=8.5, color=MPL_ORANGE, style="italic")

    # 데이터 소스
    _box(ax, (0.3, 0.35), 5.4, 0.85,
         "데이터: 도로교통법 PDF(2026.7.1 시행, 법률 제21246호) · 판례목록 CSV",
         "#EFEFEF", MPL_GRAY, 9.5, MPL_DARK)
    _box(ax, (6.1, 0.35), 5.6, 0.85,
         "가드레일: 검색 근거 밖 생성 금지 · 형량/과실 확정 금지 · 출처 ID 강제",
         "#F7E7E5", MPL_ORANGE, 9.5, "#7A2018")

    ax.set_title("도로교통법 상담 Agent — 전체 워크플로", fontsize=15, weight="bold", color=MPL_DARK, pad=12)
    fig.tight_layout()
    fig.savefig(WORKFLOW_IMG, dpi=170, bbox_inches="tight")
    plt.close(fig)


def build_langgraph_diagram():
    """LangGraph node/edge/state 다이어그램."""
    fig, ax = plt.subplots(figsize=(12, 6.2))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6.2)
    ax.axis("off")

    cx = 3.3
    # 노드(세로 배치)
    _box(ax, (cx, 5.35), 2.4, 0.7, "START", "#CFD8E3", MPL_DARK, 11, MPL_DARK)
    n_analyze = _box(ax, (cx, 4.25), 2.4, 0.8, "analyze\n쟁점 구조화", MPL_ACCENT, MPL_ACCENT, 11)
    n_retrieve = _box(ax, (cx, 3.05), 2.4, 0.8, "retrieve\n이중 하이브리드 검색", MPL_GREEN, MPL_GREEN, 11)
    n_answer = _box(ax, (cx, 1.85), 2.4, 0.8, "answer\n근거 제한 답변", MPL_DARK, MPL_DARK, 11)
    _box(ax, (cx, 0.75), 2.4, 0.7, "END", "#CFD8E3", MPL_DARK, 11, MPL_DARK)

    _arrow(ax, (cx + 1.2, 5.35), (cx + 1.2, 5.05))
    _arrow(ax, (cx + 1.2, 4.25), (cx + 1.2, 3.85))
    _arrow(ax, (cx + 1.2, 3.05), (cx + 1.2, 2.65))
    _arrow(ax, (cx + 1.2, 1.85), (cx + 1.2, 1.45))

    # 엣지 라벨 (노드 좌측, 순차 edge)
    for y in (5.2, 4.05, 2.85, 1.65):
        ax.text(cx - 0.15, y, "edge", fontsize=8, color=MPL_GRAY, va="center", ha="right")

    # State 패널
    sx = 7.2
    _box(ax, (sx, 3.7), 4.4, 2.2,
         "AgentState (TypedDict)\n\n"
         "· messages: add_messages 누적\n"
         "· analysis: 쟁점 분석 결과(dict)\n"
         "· law_hits: 법령 근거 문서\n"
         "· case_hits: 판례 근거 문서\n"
         "· context: 인용용 근거 텍스트",
         "#F2F5F9", MPL_ACCENT, 10.5, MPL_DARK)
    ax.text(sx + 2.2, 5.62, "공유 상태(State)", fontsize=9, color=MPL_ACCENT, ha="center", weight="bold")

    # State 읽기/쓰기 연결
    _arrow(ax, (cx + 2.4, 4.65), (sx, 5.0), MPL_ACCENT, ls=":")
    _arrow(ax, (cx + 2.4, 3.45), (sx, 4.6), MPL_GREEN, ls=":")
    _arrow(ax, (cx + 2.4, 2.25), (sx, 4.1), MPL_DARK, ls=":")

    # Checkpointer
    _box(ax, (sx, 1.6), 4.4, 1.5,
         "Checkpointer: MemorySaver\n\n"
         "thread_id 별 상태 저장 →\n후속 질문에서 대화 맥락 유지",
         "#FBEEDD", MPL_ORANGE, 10.5, "#7A4B12")
    _arrow(ax, (n_answer[0] + 0.0, 1.85), (sx, 2.1), MPL_ORANGE, ls="--")

    ax.set_title("LangGraph 구성 — Node · Edge · State", fontsize=15, weight="bold", color=MPL_DARK, pad=12)
    fig.tight_layout()
    fig.savefig(LANGGRAPH_IMG, dpi=170, bbox_inches="tight")
    plt.close(fig)


def load_eval_stats(csv_path: Path = EVAL_CSV) -> dict:
    """Gold Test Set 평가 결과 CSV를 읽어 발표용 요약 통계를 계산합니다."""
    df = pd.read_csv(csv_path, encoding="utf-8-sig")
    df["difficulty"] = df["difficulty"].astype(str).str.lower().str.strip()

    by_difficulty = df.groupby("difficulty").agg(
        문항수=("id", "count"),
        평균점수=("total_score", "mean"),
        통과율=("passed", "mean"),
    ).reindex(["low", "medium", "hard"]).dropna(how="all")

    component_avg = {
        "결론 일치\n(0~4)": df["conclusion_score"].mean(),
        "조문 적용\n(0~2)": df["article_score"].mean(),
        "조건부 판단\n(0~2)": df["abstention_score"].mean(),
        "금지 주장 회피\n(0~2)": df["forbidden_claim_score"].mean(),
    }

    return {
        "df": df,
        "n_total": len(df),
        "avg_score": df["total_score"].mean(),
        "pass_rate": df["passed"].mean(),
        "avg_seconds": df["response_seconds"].mean(),
        "min_seconds": df["response_seconds"].min(),
        "max_seconds": df["response_seconds"].max(),
        "n_errors": int((df["error"].fillna("").astype(str).str.strip() != "").sum()),
        "n_adopted": int((df["adopted_forbidden_claims"].fillna("").astype(str).str.strip() != "").sum()),
        "by_difficulty": by_difficulty,
        "component_avg": component_avg,
    }


def build_eval_summary_chart(stats: dict, out_path: Path = EVAL_SUMMARY_IMG, title_prefix: str = "Gold Test Set 평가 결과"):
    """난이도별 평균점수·통과율 + 항목별 평균점수를 한 이미지에 정리합니다."""
    by_diff = stats["by_difficulty"]
    component_avg = stats["component_avg"]

    fig, axes = plt.subplots(1, 2, figsize=(12, 4.6))

    # (a) 난이도별 평균 점수 + 통과율
    ax = axes[0]
    x = range(len(by_diff))
    bars = ax.bar(x, by_diff["평균점수"], color=[MPL_ACCENT, MPL_ORANGE, MPL_DARK][: len(by_diff)], width=0.5)
    ax.set_xticks(list(x))
    ax.set_xticklabels([f"{lvl}\n(n={int(n)})" for lvl, n in zip(by_diff.index, by_diff["문항수"])])
    ax.set_ylim(0, 10.5)
    ax.set_ylabel("평균 총점 (10점 만점)")
    ax.set_title("난이도별 평균 점수 · 통과율", fontsize=12, weight="bold", color=MPL_DARK)
    for i, (score, rate) in enumerate(zip(by_diff["평균점수"], by_diff["통과율"])):
        ax.text(i, score + 0.25, f"{score:.2f}점", ha="center", fontsize=10, weight="bold", color=MPL_DARK)
        ax.text(i, score - 0.9, f"통과율 {rate * 100:.0f}%", ha="center", fontsize=9, color="white", weight="bold")
    ax.axhline(7, color=MPL_GRAY, linestyle="--", linewidth=1)
    ax.text(len(by_diff) - 0.5, 7.15, "통과 기준 7점", fontsize=8, color=MPL_GRAY, ha="right")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)

    # (b) 채점 항목별 평균 점수
    ax2 = axes[1]
    labels = list(component_avg.keys())
    values = list(component_avg.values())
    maxes = [4, 2, 2, 2]
    y = range(len(labels))
    ax2.barh(y, values, color=MPL_GREEN, height=0.55)
    ax2.barh(y, maxes, color="#E4E4E4", height=0.55, zorder=0)
    ax2.barh(y, values, color=MPL_GREEN, height=0.55, zorder=1)
    ax2.set_yticks(list(y))
    ax2.set_yticklabels(labels, fontsize=10)
    ax2.invert_yaxis()
    ax2.set_xlim(0, 4.3)
    ax2.set_title("채점 항목별 평균 점수", fontsize=12, weight="bold", color=MPL_DARK)
    for i, (v, m) in enumerate(zip(values, maxes)):
        ax2.text(v + 0.08, i, f"{v:.2f}/{m}", va="center", fontsize=9.5, color=MPL_DARK, weight="bold")
    for spine in ("top", "right"):
        ax2.spines[spine].set_visible(False)

    fig.suptitle(
        f"{title_prefix} (전체 {stats['n_total']}문항, 평균 {stats['avg_score']:.2f}점 · 통과율 {stats['pass_rate'] * 100:.0f}%)",
        fontsize=13, weight="bold", color=MPL_DARK, y=1.03,
    )
    fig.tight_layout()
    fig.savefig(out_path, dpi=170, bbox_inches="tight")
    plt.close(fig)


# ── PPT 유틸 ────────────────────────────────────────────────
def _title_only(prs):
    return prs.slides.add_slide(prs.slide_layouts[5])


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _style_title(slide, text, color=DARK, size=28):
    slide.shapes.title.text = text
    p = slide.shapes.title.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color


def _accent_bar(slide, prs, color=ACCENT):
    """제목 아래 얇은 강조 바."""
    bar = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.35), Inches(9.0), Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _bullets(slide, items, left=0.7, top=1.7, width=8.9, height=5.2, size=17, gap=6):
    tf = _textbox(slide, left, top, width, height)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "" if level > 0 else "• "
        p.text = ("    " * level) + prefix + text
        p.font.size = Pt(size - level * 2)
        p.font.color.rgb = GRAY if level > 0 else DARK
        p.space_after = Pt(gap)


def _two_col(slide, left_title, left_items, right_title, right_items,
             left_color=ACCENT2, right_color=WARN):
    # 왼쪽 카드
    for x, ttl, items, col in (
        (0.6, left_title, left_items, left_color),
        (5.15, right_title, right_items, right_color),
    ):
        header = slide.shapes.add_shape(1, Inches(x), Inches(1.75), Inches(4.25), Inches(0.55))
        header.fill.solid()
        header.fill.fore_color.rgb = col
        header.line.fill.background()
        htf = header.text_frame
        htf.text = ttl
        hp = htf.paragraphs[0]
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

        tf = _textbox(slide, x + 0.05, 2.45, 4.15, 4.4)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)


def _section(prs, num, title, subtitle=""):
    slide = _title_only(prs)
    _set_bg(slide, DARK)
    slide.shapes.title.text = f"Chapter {num}"
    tp = slide.shapes.title.text_frame.paragraphs[0]
    tp.font.size = Pt(20)
    tp.font.color.rgb = ACCENT
    tp.font.bold = True

    tf = _textbox(slide, 0.8, 2.9, 8.5, 2.0)
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(0xB9, 0xC7, 0xD8)
    return slide


def _image_slide(prs, title, image_path, notes):
    slide = _title_only(prs)
    _style_title(slide, title, size=24)
    _accent_bar(slide, prs)
    if image_path.exists():
        # 슬라이드 폭(10인치)에 맞춰 배치
        slide.shapes.add_picture(str(image_path), Inches(0.55), Inches(1.55), width=Inches(8.9))
    tf = _textbox(slide, 0.7, 6.35, 8.9, 1.0)
    for i, note in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + note
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
    return slide


def build_ppt(stats: dict, wo_bm25_stats: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 표지 ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "도로교통법 전문 상담 AI Agent"
    tp = slide.shapes.title.text_frame.paragraphs[0]
    tp.font.size = Pt(34)
    tp.font.bold = True
    tp.font.color.rgb = DARK
    sub = slide.placeholders[1]
    sub.text = (
        "LangGraph 기반 하이브리드 RAG 법률 상담 챗봇\n"
        "개발자 소개 자료  ·  doitmyself/traffic_law_z.ipynb  ·  2026.07"
    )
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = ACCENT

    # ── 목차 ──
    slide = _title_only(prs)
    _style_title(slide, "목차")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "Chapter 1. 개요 — 목적 · 가능/불가능 · 신뢰도",
        "Chapter 2. 사용한 기술 — 모델 · 프레임워크 · Tool · 전략 · 프롬프트",
        "Chapter 3. Workflow — 상담 흐름도 · LangGraph 구조",
        "Chapter 4. 검증 — Gold Test Set · 자동 채점 · Ablation 비교",
    ], size=20, gap=14)

    # ══ Chapter 1 ══
    _section(prs, 1, "개요", "무엇을, 왜, 얼마나 믿을 수 있는가")

    # 1-1 목적
    slide = _title_only(prs)
    _style_title(slide, "1-1. 이 챗봇의 목적")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "도로교통법 관련 사고·벌금·법 위반 질문에 근거 기반으로 법적 판단을 보조",
        ("현행 도로교통법(2026.7.1 시행, 법률 제21246호) 조문과 대법원 판례 요약을 검색", 1),
        ("검색된 근거만 인용해 답변 — 환각(hallucination) 최소화", 1),
        "일반 사용자가 이해할 수 있는 상담형 답변 제공",
        ("'잠정 판단 → 적용 근거 → 추가 필요 사실 → 권장 조치' 구조로 정리", 1),
        "대상: 법률 비전문가 사용자 + Agentic RAG 패턴을 학습하려는 개발자",
    ], size=17)

    # 1-2 할 수 있는 것 / 없는 것
    slide = _blank(prs)
    tf = _textbox(slide, 0.6, 0.55, 8.8, 0.8)
    tf.text = "1-2. 할 수 있는 것 vs 할 수 없는 것"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    _two_col(
        slide,
        "할 수 있는 것", [
            "관련 조문·판례 검색 및 인용",
            "위반 여부의 조건부·잠정 판단",
            "누락된 핵심 사실 되묻기",
            "사고 직후 초동 조치 안내",
            "출처 ID(L·C)로 근거 추적",
            "후속 질문의 문맥 유지",
        ],
        "할 수 없는 것", [
            "최종 형량·벌금액 확정",
            "과실비율·유무죄 확정 판단",
            "변호사 법률 자문 대체",
            "시행령·시행규칙 등 미수록 법령 단정",
            "제공 자료 밖 사실 생성",
            "실시간 최신 개정 자동 반영",
        ],
    )

    # 1-3 신뢰도
    slide = _title_only(prs)
    _style_title(slide, "1-3. 결과의 신뢰도")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "신뢰를 높이는 장치",
        ("검색 근거 제한: 제공된 법령·판례 밖 내용 생성 금지", 1),
        ("출처 강제 인용: 핵심 판단마다 [L1]/[C1] 근거 ID 표기", 1),
        ("기준일 명시: 2026.7.1 시행 현행법 기준임을 항상 고지", 1),
        ("Gold Test Set으로 정확도·안전성 정량 검증 (Chapter 4)", 1),
        "남는 한계",
        ("판례는 '요약본'이라 사건별 세부 사실관계는 제한적", 1),
        ("시행령·특별법 미수록 → 실제 처분과 차이 가능", 1),
        ("최종 결과는 반드시 전문가·관할기관 확인 필요", 1),
    ], size=16)

    # ══ Chapter 2 ══
    _section(prs, 2, "사용한 기술들", "모델 · 프레임워크 · Tool · 전략 · 프롬프트")

    # 2-1 모델 & 프레임워크
    slide = _title_only(prs)
    _style_title(slide, "2-1. LLM 모델 & 프레임워크")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "LLM",
        ("Chat: gpt-4.1-mini (temperature=0) — 분석·답변·자동 채점", 1),
        ("Embedding: text-embedding-3-small — 법령·판례 벡터화", 1),
        ("모델 교체 용이: llm / embeddings 객체만 호환 구현으로 대체 가능", 1),
        "프레임워크 (워크프레임)",
        ("LangGraph: StateGraph 기반 상담 워크플로 오케스트레이션", 1),
        ("LangChain: 문서·메시지·구조화 출력(with_structured_output)", 1),
        ("Chroma: 영속 벡터 DB · pypdf: PDF 파싱 · pandas: 판례 CSV", 1),
        ("Gradio: 웹 채팅 UI · Pydantic: 상태·스키마 검증", 1),
    ], size=16)

    # 2-2 Tool
    slide = _title_only(prs)
    _style_title(slide, "2-2. Tool (검색 도구)")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "벡터 검색 (Semantic) — Chroma similarity_search",
        ("의미 기반 유사도로 문맥이 비슷한 조문·판례 검색", 1),
        "BM25 키워드 검색 (rank-bm25 / BM25Okapi)",
        ("'제44조', 정확한 법률 용어 등 키워드 정확도 보완", 1),
        "RRF (Reciprocal Rank Fusion)",
        ("두 검색 순위를 score = weight/(60+rank)로 융합", 1),
        ("키워드 가중치 1.15로 법률 용어 매칭 우대", 1),
        "이중 인덱스: 법령 전용 / 판례 전용 컬렉션 분리 검색",
        "구조화 추출 Tool: with_structured_output으로 쟁점 JSON 생성",
    ], size=16)

    # 2-3 전략
    slide = _title_only(prs)
    _style_title(slide, "2-3. 전략 (Chunking · Retriever 고도화)")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "Chunking",
        ("법령: 페이지 단위 정제 후 문장 경계 우선 분할(≈1100자, overlap 160)", 1),
        ("조문 번호(제N조/제N조의M) 메타데이터 자동 태깅 → 출처 추적", 1),
        ("판례: 1건=1문서(사건명·번호·선고일·조문·판시사항 구조화)", 1),
        "Retriever 고도화",
        ("벡터 + BM25 하이브리드 → RRF 순위 융합", 1),
        ("다중 쿼리: 분석 단계의 search_queries + 조문 후보로 재검색", 1),
        ("법령7 / 판례5건 수집 후 중복 제거", 1),
        "운영 전략",
        ("데이터 지문(fingerprint) 기반 컬렉션 캐싱 → 중복 임베딩 방지", 1),
    ], size=15)

    # 2-4 프롬프트 엔지니어링
    slide = _title_only(prs)
    _style_title(slide, "2-4. 프롬프트 엔지니어링")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "역할 분리 페르소나 (2단계)",
        ("분석 Agent: '도로교통 사건의 쟁점을 구조화하는 분석가'", 1),
        ("상담 AI: '대한민국 도로교통법 전문 상담 AI'", 1),
        "답변 가드레일 (시스템 프롬프트에 규칙 명시)",
        ("검색 근거 밖 법령·벌금·판례 생성 금지", 1),
        ("핵심 판단마다 근거 ID 인용 강제", 1),
        ("형량·과실·유무죄 확정 금지, 부족 시 조건부 답변 + 되묻기", 1),
        ("사고 직후엔 인명 구조·신고 등 초동 조치 우선 안내", 1),
        ("답변 말미에 '법률 자문이 아님' 고지 삽입", 1),
        "출력 구조 고정: 잠정 판단 → 적용 근거 → 추가 필요 사실 → 권장 조치",
    ], size=15)

    # ══ Chapter 3 ══
    _section(prs, 3, "Workflow", "상담 흐름도와 LangGraph 구조")

    _image_slide(
        prs, "3-1. 전체 상담 워크플로 흐름도", WORKFLOW_IMG,
        [
            "① 쟁점 분석: 질문에서 유형·쟁점·검색어·추정 조문·누락 사실을 구조화 추출",
            "② 이중 검색: 법령/판례를 각각 벡터+BM25로 검색 후 RRF로 융합",
            "③ 근거 제한 답변: 인용 ID 기반 구조화 답변 + 안전 고지, thread_id로 문맥 유지",
        ],
    )

    _image_slide(
        prs, "3-2. LangGraph 구조 (Node · Edge · State)", LANGGRAPH_IMG,
        [
            "Node: analyze(쟁점 구조화) → retrieve(이중 하이브리드 검색) → answer(근거 답변)",
            "Edge: START→analyze→retrieve→answer→END 순차 실행, 각 노드가 State를 갱신",
            "State: messages/analysis/law_hits/case_hits/context, MemorySaver로 세션 유지",
        ],
    )

    # ══ Chapter 4 ══
    _section(prs, 4, "검증", "Gold Test Set과 자동 채점")

    # 4-1 golden set 구조
    slide = _title_only(prs)
    _style_title(slide, "4-1. Gold Test Set 구조")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "전문가가 정답을 라벨링한 도로교통법 QA 검증셋 (10·50문항)",
        "주요 컬럼",
        ("question / difficulty(low·medium·hard) — 질문과 난이도", 1),
        ("required_articles — 반드시 근거로 삼아야 할 조문", 1),
        ("legal_issue / key_facts / missing_facts — 쟁점·핵심·누락 사실", 1),
        ("expected_conclusion — 기대되는 결론", 1),
        ("should_abstain — 단정 대신 조건부로 답해야 하는지 여부", 1),
        ("must_not_claim — 말하면 안 되는 오답·금지 주장", 1),
        ("난이도 추출 설정: even / easy(low:med:hard=5:3:2) / 단일 난이도", 1),
    ], size=15)

    # 4-2 검증 방법
    slide = _title_only(prs)
    _style_title(slide, "4-2. 검증 방법 (LLM-as-a-Judge)")
    _accent_bar(slide, prs)
    _bullets(slide, [
        "문항을 독립 대화(thread)로 순차 질의 → Agent 답변 수집",
        "LLM 심판이 Gold 기준과 대조해 항목별 채점 (총 10점)",
        ("결론 일치 0~4 · 조문 적용 0~2 · 조건부 판단 0~2 · 금지 주장 회피 0~2", 1),
        "금지 주장 회피는 채택/모호/회피 목록으로 분류 후 코드가 점수 확정",
        ("adopted→0 · ambiguous→1 · none→2 (심판 정수 편향 보정)", 1),
        "통과 기준: 총 7점 이상 + 결론 3점 이상 + 금지 주장 미채택",
        "참고 지표: 문항별 답변 생성 시간(초) 기록",
        "산출물: 문항별 점수·근거·응답시간 CSV 저장 + 난이도별 요약",
    ], size=15)

    # 4-3 실제 평가 결과
    by_diff = stats["by_difficulty"]
    diff_kr = {"low": "쉬움", "medium": "중간", "hard": "어려움"}
    diff_line = " · ".join(
        f"{diff_kr.get(lvl, lvl)} {row['평균점수']:.2f}점/통과 {row['통과율'] * 100:.0f}%"
        for lvl, row in by_diff.iterrows()
    )
    slide = _title_only(prs)
    _style_title(slide, "4-3. 실제 평가 결과 (evaluation_results_9b5c2dbe)")
    _accent_bar(slide, prs)
    if EVAL_SUMMARY_IMG.exists():
        slide.shapes.add_picture(str(EVAL_SUMMARY_IMG), Inches(0.4), Inches(1.55), width=Inches(9.2))
    tf = _textbox(slide, 0.7, 6.05, 8.9, 1.3)
    bullets = [
        f"전체 {stats['n_total']}문항 · 평균 {stats['avg_score']:.2f}/10점 · 통과율 {stats['pass_rate'] * 100:.0f}% "
        f"(오류 {stats['n_errors']}건, 평균 응답 {stats['avg_seconds']:.1f}초)",
        f"난이도별: {diff_line}",
        "주요 실패 원인: 금지 주장을 명확히 부정하지 않은 '모호한 회피'(감점 1점) · 일부 문항 근거 조문 누락",
        "금지 주장을 결론으로 채택(0점)한 사례는 0건 — 위험한 오답은 발생하지 않음",
    ]
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(4)

    # 4-4 Ablation: BM25 제거(벡터 전용) 버전 평가 결과
    if wo_bm25_stats is not None:
        wdiff = wo_bm25_stats["by_difficulty"]
        wdiff_line = " · ".join(
            f"{diff_kr.get(lvl, lvl)} {row['평균점수']:.2f}점/통과 {row['통과율'] * 100:.0f}%"
            for lvl, row in wdiff.iterrows()
        )
        delta_score = wo_bm25_stats["avg_score"] - stats["avg_score"]
        delta_pass = (wo_bm25_stats["pass_rate"] - stats["pass_rate"]) * 100
        delta_time = wo_bm25_stats["avg_seconds"] - stats["avg_seconds"]

        slide = _title_only(prs)
        _style_title(slide, "4-4. Ablation — BM25 제거(벡터 전용) 비교 (evaluation_results_a5cb495c)")
        _accent_bar(slide, prs)
        if EVAL_WO_BM25_SUMMARY_IMG.exists():
            slide.shapes.add_picture(str(EVAL_WO_BM25_SUMMARY_IMG), Inches(0.4), Inches(1.55), width=Inches(9.2))
        tf = _textbox(slide, 0.7, 6.05, 8.9, 1.3)
        bullets2 = [
            f"전체 {wo_bm25_stats['n_total']}문항 · 평균 {wo_bm25_stats['avg_score']:.2f}/10점 · 통과율 {wo_bm25_stats['pass_rate'] * 100:.0f}% "
            f"(오류 {wo_bm25_stats['n_errors']}건, 평균 응답 {wo_bm25_stats['avg_seconds']:.1f}초)",
            f"난이도별: {wdiff_line}",
            f"하이브리드(BM25+벡터) 대비: 평균점수 {delta_score:+.2f}점 · 통과율 {delta_pass:+.0f}%p · 응답시간 {delta_time:+.1f}초",
            "BM25 제거 시 금지 주장을 결론으로 채택(0점)한 사례가 발생 — 키워드 검색이 정확한 조문·수치 매칭에 기여",
        ]
        for i, text in enumerate(bullets2):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + text
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(4)

    # ── 마무리 ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "감사합니다"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK
    sub = slide.placeholders[1]
    sub.text = (
        "코드: doitmyself/traffic_law_z.ipynb\n"
        "데이터: 도로교통법 PDF · 판례목록 CSV · Gold Test Set\n"
        "⚠️ 본 챗봇은 정보 제공용이며 법률 자문을 대체하지 않습니다."
    )
    sub.text_frame.paragraphs[0].font.size = Pt(16)
    sub.text_frame.paragraphs[0].font.color.rgb = ACCENT

    prs.save(str(OUT_PPTX))
    print(f"saved: {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_workflow_diagram()
    print(f"saved: {WORKFLOW_IMG}")
    build_langgraph_diagram()
    print(f"saved: {LANGGRAPH_IMG}")

    eval_stats = load_eval_stats(EVAL_CSV)
    build_eval_summary_chart(eval_stats, EVAL_SUMMARY_IMG, "Gold Test Set 평가 결과 (하이브리드 RAG)")
    print(f"saved: {EVAL_SUMMARY_IMG}")

    wo_bm25_stats = None
    if EVAL_WO_BM25_CSV.exists():
        wo_bm25_stats = load_eval_stats(EVAL_WO_BM25_CSV)
        build_eval_summary_chart(wo_bm25_stats, EVAL_WO_BM25_SUMMARY_IMG, "Gold Test Set 평가 결과 (벡터 전용 RAG)")
        print(f"saved: {EVAL_WO_BM25_SUMMARY_IMG}")
    else:
        print(f"skip: {EVAL_WO_BM25_CSV} not found")

    build_ppt(eval_stats, wo_bm25_stats)
