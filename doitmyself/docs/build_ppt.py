"""simple_lang1_overview.md 기반 PPT 생성."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).parent
OUT = DOCS / "simple_lang1_overview.pptx"

# 색상 테마
DARK = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_title(slide, title: str, subtitle: str = ""):
    """슬라이드 제목·부제 설정."""
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullets(slide, items: list[str], font_size: int = 18):
    """본문 불릿 리스트 추가."""
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(items):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = GRAY


def add_title_bar(slide, prs):
    """제목 슬라이드 스타일."""
    title = slide.shapes.title
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK
    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = ACCENT


def add_section_slide(prs, title: str):
    """섹션 구분 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = WHITE
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
    return slide


def add_image_slide(prs, title: str, image_path: Path, caption: str = ""):
    """이미지 중심 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK
    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5), Inches(1.3),
            width=Inches(9),
        )
    if caption:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = box.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = GRAY
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "도로교통법 Agent PoC"
    slide.placeholders[1].text = (
        "LangGraph 기반 Agentic RAG 소개\n"
        "doitmyself/simple_lang1.ipynb\n"
        "2026.07"
    )
    add_title_bar(slide, prs)

    # 2. 목차
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "목차", "")
    add_bullets(slide, [
        "1. 애플리케이션 개요",
        "2. 기술 스택 · 모델 · Tool",
        "3. Agentic RAG 및 설계 패턴",
        "4. LangGraph 흐름도",
        "5. 세션 메모리 · 문맥 관리",
        "6. 실행 방법 · 확장 포인트",
    ], font_size=22)

    # 3. 목적
    add_section_slide(prs, "1. 애플리케이션 개요")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "목적", "도로교통법 Q&A Agent PoC")
    add_bullets(slide, [
        "법률 조문 · 판례 · 웹 사례를 근거로 상황을 해석해 답변",
        "LangChain / LangGraph 학습 패턴을 실전형 예제로 통합",
        "day4 RAG + day15 Graph + day16 Tool/Chroma 패턴 적용",
        "대상: AI Agent 개발자 · PoC 데모 · FastAPI/Slack 연동 기반",
    ])

    # 4. 주요 기능
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "주요 기능", "")
    add_bullets(slide, [
        "법률 Q&A — PDF RAG 검색 후 일반인용 해설",
        "판례 기반 해석 — 사고·과실 판단 질문 (CASE)",
        "웹 검색 폴백 — DuckDuckGo 유사 사례 + 링크",
        "triage — 무관 질문 거절 / clarify — 역질문 (최대 2회)",
        "[근거] 표시 — 법률 PDF / 판례 / 웹 링크 출처 명시",
        "멀티 세션 (user1~3) + Rolling Summary 장기 문맥",
    ], font_size=17)

    # 5. 데이터 소스
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "데이터 소스", "")
    add_bullets(slide, [
        "simple_lang1_dataset/도로교통법_법률_제21246호.pdf → 법률 RAG",
        "simple_lang1_dataset/판례목록.txt → 판례 임베딩 (줄 단위)",
        "chroma_traffic_law/ — 법률 벡터 DB",
        "chroma_case_list/ — 판례 벡터 DB",
        "REBUILD_CHROMA=True 시 벡터 DB 재구축",
    ])

    # 6. 기술 스택
    add_section_slide(prs, "2. 기술 스택 · 모델 · Tool")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "프레임워크 & 모델", "")
    add_bullets(slide, [
        "오케스트레이션: LangGraph (StateGraph, Conditional Edge)",
        "LLM: gpt-4o-mini (temperature=0) — triage·분류·생성",
        "Embedding: text-embedding-3-small — PDF + 판례목록",
        "벡터 DB: Chroma | 문서: PyPDFLoader + TextSplitter",
        "상태: Pydantic AgentState | 웹: ddgs (DuckDuckGo)",
    ], font_size=17)

    # 7. Tool
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Tool (검색 도구)", "Graph 노드 = Tool 실행 단위")
    add_bullets(slide, [
        "법률 RAG — vectorstore.similarity_search(k=4)",
        "판례 검색 — similarity_search_with_score + threshold 0.40",
        "웹 검색 — web_search_cases() → DDGS().text(max=3)",
        "튜닝: MAX_CLARIFY=2, RECENT=5턴, SUMMARY=6턴 이후 압축",
    ])

    # 8. Agentic RAG
    add_section_slide(prs, "3. Agentic RAG 및 설계 패턴")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Agentic RAG란?", "단순 RAG vs Multi-step RAG")
    add_bullets(slide, [
        "단순 RAG: 질문 → 검색 → 생성",
        "Agentic RAG: LLM이 라우팅·분류·도구 선택을 결정",
        "질문 → [triage] → [retrieve] → [classify]",
        "  → interpret(LAW) 또는 case_lookup(CASE) → web 폴백",
        "검색 전·후에 LLM 판단 노드가 끼어 있는 구조",
    ], font_size=17)

    # 9. 노드 역할
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "LangGraph 노드 (9개)", "")
    add_bullets(slide, [
        "triage — retrieve / clarify / reject 라우팅 (LLM)",
        "retrieve — 법률 PDF RAG | classify — LAW vs CASE",
        "interpret — 법률 해설 | case_lookup — 판례 검색",
        "case_answer — 판례 해설 | web_search — 웹 폴백",
        "clarify — 역질문 | reject — 거절 메시지",
        "AgentState: question, history_text, law_docs, answer, source ...",
    ], font_size=16)

    # 10. 문맥 & API
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "문맥 관리 & 세션 API", "")
    add_bullets(slide, [
        "_context_block() — 모든 답변·분류 노드에 문맥 주입",
        "_search_query() — 후속 질문 RAG/판례/웹 검색 쿼리",
        "Rolling Summary — 6턴 초과 시 LLM 요약, 최근 5턴 원문",
        "run_turn(session_id, question) → dict 반환",
        "  (Slack/FastAPI 연동 준비: answer, source, context_summary)",
        "명령어: 사용자변경 | 세션비우기 | 종료",
    ], font_size=16)

    # 11. 흐름도
    add_section_slide(prs, "4. LangGraph 흐름도")
    add_image_slide(
        prs,
        "Agent 전체 흐름",
        DOCS / "simple_lang1_agent_flow.png",
        "START → triage → retrieve/classify → interpret 또는 case_lookup → END",
    )

    # 12. 세션 메모리
    add_section_slide(prs, "5. 세션 메모리 · 문맥 관리")
    add_image_slide(
        prs,
        "세션 메모리 구조",
        DOCS / "simple_lang1_session_memory.png",
        "prepare_graph_input → app.invoke → save_session_turn",
    )

    # 13. 실행 & 테스트
    add_section_slide(prs, "6. 실행 · 확장")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Quick Start & 테스트", "")
    add_bullets(slide, [
        ".env에 OPENAI_API_KEY 설정",
        "simple_lang1.ipynb Step 0~12 순서 실행",
        "Step 12: 세션 ID 입력 후 대화",
        "테스트 A: 무관 질문 → reject",
        "테스트 B: 처벌 기준 → LAW → interpret",
        "테스트 C: 사고 과실 → CASE → case_answer",
        "테스트 D~F: clarify 재시도 → reject",
    ], font_size=16)

    # 14. 확장 포인트
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "확장 포인트", "")
    add_bullets(slide, [
        "API 서버 — run_turn() + FastAPI/Slack 웹훅",
        "Checkpointer — MemorySaver / Redis 영속화",
        "Tool 표준화 — @tool + ToolNode",
        "판례 본문 RAG — 판결문 PDF 추가",
        "평가 — triage/classify 벤치마크",
    ])

    # 15. 마무리
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "감사합니다"
    slide.placeholders[1].text = (
        "doitmyself/simple_lang1.ipynb\n"
        "doitmyself/docs/simple_lang1_overview.md"
    )
    add_title_bar(slide, prs)

    prs.save(str(OUT))
    print(f"saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
